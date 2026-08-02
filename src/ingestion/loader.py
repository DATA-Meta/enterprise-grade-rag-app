"""
Document ingestion: load raw files from a folder and split them into chunks.

Supported formats: .txt, .pdf, .docx, .pptx, .html
"""

from dataclasses import dataclass
from pathlib import Path

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup


@dataclass
class Document:
    """One loaded file, before chunking."""
    source: str   # filename, so we can trace a chunk back to its origin
    text: str     # full raw text of the file


@dataclass
class Chunk:
    """One chunk of a document, ready to be embedded."""
    source: str
    chunk_id: int
    text: str


def load_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def load_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)


def load_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def load_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    text_runs = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text.strip():
                        text_runs.append(text)
    return "\n".join(text_runs)


def load_html(path: Path) -> str:
    raw_html = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(raw_html, "html.parser")
    # Remove script/style tags so we don't ingest code as if it were content
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n")


LOADERS = {
    ".txt": load_txt,
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".html": load_html,
    ".htm": load_html,
}


def load_documents(data_dir: str) -> list[Document]:
    """
    Read every supported file in `data_dir` and return a list of Document objects.
    """
    data_path = Path(data_dir)
    if not data_path.exists():
        raise FileNotFoundError(f"Data directory not found: {data_dir}")

    documents: list[Document] = []

    for file_path in data_path.iterdir():
        if not file_path.is_file():
            continue

        suffix = file_path.suffix.lower()
        loader_fn = LOADERS.get(suffix)

        if loader_fn is None:
            # Unsupported file type — skip it rather than crash.
            continue

        try:
            text = loader_fn(file_path)
        except Exception as e:
            print(f"Warning: failed to load {file_path.name}: {e}")
            continue

        if text.strip():
            documents.append(Document(source=file_path.name, text=text))

    return documents


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 150) -> list[str]:
    """
    Split text into overlapping chunks of roughly `chunk_size` characters.
    """
    if chunk_size <= overlap:
        raise ValueError("chunk_size must be greater than overlap")

    chunks = []
    start = 0
    text_length = len(text)

    while start < text_length:
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap

    return chunks


def chunk_documents(documents: list[Document], chunk_size: int = 1000, overlap: int = 150) -> list[Chunk]:
    """
    Turn a list of Documents into a flat list of Chunks, ready for embedding.
    """
    all_chunks: list[Chunk] = []

    for doc in documents:
        pieces = chunk_text(doc.text, chunk_size=chunk_size, overlap=overlap)
        for i, piece in enumerate(pieces):
            all_chunks.append(Chunk(source=doc.source, chunk_id=i, text=piece))

    return all_chunks


if __name__ == "__main__":
    docs = load_documents("DATA")
    print(f"Loaded {len(docs)} document(s).")
    for doc in docs:
        print(f"  - {doc.source} ({len(doc.text)} characters)")

    chunks = chunk_documents(docs)
    print(f"Produced {len(chunks)} chunk(s).")

    if chunks:
        print("\n--- First chunk preview ---")
        print(f"Source: {chunks[0].source} | Chunk ID: {chunks[0].chunk_id}")
        print(chunks[0].text[:300])